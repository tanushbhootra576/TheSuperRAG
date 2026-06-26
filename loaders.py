import os
import re
import csv
from pathlib import Path
from typing import List, Union
from urllib.parse import urlparse
from datetime import datetime

from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader

class DocumentLoader:
    @classmethod
    def _metadata(cls, source: str, file_type: str, page_number: Union[int, str] = "N/A", title: str = "") -> dict:
        return {
            "source": str(source),
            "file_type": file_type,
            "page_number": page_number,
            "created_at": datetime.utcnow().isoformat(),
            "title": title or str(source),
        }

    @classmethod
    def load(cls, source: Union[str, Path]) -> List[Document]:
        source_str = str(source)
        
        # 1. URL detection
        if source_str.startswith("http://") or source_str.startswith("https://"):
            if "youtube.com" in source_str or "youtu.be" in source_str:
                return cls._load_youtube(source_str)
            else:
                return cls._load_web(source_str)
                
        # 2. Local File detection
        path = Path(source)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {source}")
            
        ext = path.suffix.lower()
        if ext == ".pdf":
            return cls._load_pdf(path)
        elif ext == ".docx":
            return cls._load_docx(path)
        elif ext in [".txt", ".md"]:
            return cls._load_text(path)
        elif ext == ".csv":
            return cls._load_csv(path)
        elif ext == ".xlsx":
            return cls._load_xlsx(path)
        elif ext == ".pptx":
            return cls._load_pptx(path)
        elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
            return cls._load_image(path)
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

    @classmethod
    def _load_pdf(cls, path: Path) -> List[Document]:
        import fitz
        try:
            import pytesseract
            has_tesseract = True
        except ImportError:
            has_tesseract = False

        docs = []
        doc = fitz.open(str(path))
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text().strip()

            # Fallback to OCR if page has no extractable text (scanned PDF)
            if not text and has_tesseract:
                try:
                    pix = page.get_pixmap(dpi=150)
                    from PIL import Image
                    import io
                    img = Image.open(io.BytesIO(pix.tobytes("png")))
                    text = pytesseract.image_to_string(img).strip()
                except Exception as e:
                    print(f"  [X] OCR failed for {path.name} page {page_num + 1}: {e}")

            if text:
                docs.append(Document(
                    page_content=text,
                    metadata=cls._metadata(path.name, "pdf", page_num + 1, path.name)
                ))
        doc.close()
        return docs

    @classmethod
    def _load_docx(cls, path: Path) -> List[Document]:
        import docx
        doc = docx.Document(str(path))
        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        return [Document(
            page_content=text,
            metadata=cls._metadata(path.name, "docx", 1, path.name)
        )]

    @classmethod
    def _load_text(cls, path: Path) -> List[Document]:
        with open(path, "r", encoding="utf-8") as f:
            text = f.read()
        return [Document(
            page_content=text,
            metadata=cls._metadata(path.name, path.suffix.lower()[1:], 1, path.name)
        )]

    @classmethod
    def _load_csv(cls, path: Path) -> List[Document]:
        docs = []
        with open(path, "r", encoding="utf-8") as f:
            reader = csv.DictReader(f)
            for i, row in enumerate(reader, start=1):
                content = "\n".join([f"{k}: {v}" for k, v in row.items() if v])
                if content.strip():
                    docs.append(Document(
                        page_content=content,
                        metadata=cls._metadata(path.name, "csv", i, path.name)
                    ))
        return docs

    @classmethod
    def _load_xlsx(cls, path: Path) -> List[Document]:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True)
        docs = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            headers = [cell.value for cell in ws[1]]
            for i, row in enumerate(ws.iter_rows(min_row=2, values_only=True), start=2):
                if any(row):
                    content = "\n".join([f"{headers[j] if j < len(headers) else f'Col{j}'}: {val}" for j, val in enumerate(row) if val is not None])
                    docs.append(Document(
                        page_content=content,
                        metadata=cls._metadata(path.name, "xlsx", f"{sheet_name} R{i}", path.name)
                    ))
        return docs

    @classmethod
    def _load_pptx(cls, path: Path) -> List[Document]:
        from pptx import Presentation
        prs = Presentation(str(path))
        docs = []
        for i, slide in enumerate(prs.slides, start=1):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            content = "\n".join(text).strip()
            if content:
                docs.append(Document(
                    page_content=content,
                    metadata=cls._metadata(path.name, "pptx", i, path.name)
                ))
        return docs

    @classmethod
    def _load_youtube(cls, url: str) -> List[Document]:
        from youtube_transcript_api import YouTubeTranscriptApi
        video_id = None
        if "youtu.be/" in url:
            video_id = url.split("youtu.be/")[1].split("?")[0]
        elif "youtube.com/watch" in url:
            from urllib.parse import parse_qs
            parsed_url = urlparse(url)
            video_id = parse_qs(parsed_url.query).get("v", [None])[0]
            
        if not video_id:
            raise ValueError(f"Could not extract video ID from YouTube URL: {url}")
            
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        text = " ".join([t['text'] for t in transcript])
        return [Document(
            page_content=text,
            metadata=cls._metadata(url, "youtube", "video", url)
        )]

    @classmethod
    def _load_web(cls, url: str) -> List[Document]:
        import httpx
        from bs4 import BeautifulSoup
        resp = httpx.get(url, timeout=10.0, follow_redirects=True)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.extract()
            
        text = soup.get_text(separator="\n")
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = "\n".join(chunk for chunk in chunks if chunk)
        
        title = soup.title.string if soup.title else url
        return [Document(
            page_content=text,
            metadata=cls._metadata(url, "web", "page", title)
        )]

    @classmethod
    def _load_image(cls, path: Path) -> List[Document]:
        import pytesseract
        from PIL import Image
        img = Image.open(str(path))
        text = pytesseract.image_to_string(img)
        if not text.strip():
            text = "[No text found in image]"
        return [Document(
            page_content=text,
            metadata=cls._metadata(path.name, "image", 1, path.name)
        )]
